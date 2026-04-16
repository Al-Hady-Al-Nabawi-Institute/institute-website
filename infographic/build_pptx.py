#!/usr/bin/env python3
"""Build a PowerPoint deck matching presentation-draft.html."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

TEAL      = RGBColor(0x0d, 0x6b, 0x63)
TEAL_DARK = RGBColor(0x08, 0x4d, 0x47)
TEAL_XDARK= RGBColor(0x05, 0x38, 0x33)
GOLD      = RGBColor(0xe0, 0xb3, 0x33)
CREAM     = RGBColor(0xf7, 0xf1, 0xe3)
INK       = RGBColor(0x1a, 0x2a, 0x27)
GREY_TXT  = RGBColor(0x2a, 0x38, 0x36)
GREY_SUB  = RGBColor(0x4a, 0x58, 0x56)
WHITE     = RGBColor(0xff, 0xff, 0xff)

FONT_DISPLAY = "Noto Naskh Arabic"
FONT_BODY    = "IBM Plex Sans Arabic"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def set_rtl(paragraph):
    pPr = paragraph._pPr if paragraph._pPr is not None else paragraph._p.get_or_add_pPr()
    pPr.set('rtl', '1')


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             font=FONT_BODY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP, rtl=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if rtl:
            set_rtl(p)
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_brand_strip(slide):
    add_rect(slide, 0, SH - Emu(60000), SW, Emu(60000), GOLD)


def add_slide_number(slide, n, light=False):
    col = RGBColor(0xbb, 0xbb, 0xbb) if light else RGBColor(0x8a, 0x8a, 0x8a)
    add_text(slide, Inches(0.3), SH - Inches(0.5), Inches(1), Inches(0.3),
             f"{n:02d}", size=11, color=col, align=PP_ALIGN.LEFT, rtl=False)


def add_title(slide, text, subtitle=None):
    # right gold bar
    add_rect(slide, SW - Inches(0.55), Inches(0.55), Inches(0.08), Inches(0.9), GOLD)
    add_text(slide, Inches(0.5), Inches(0.5), SW - Inches(1.2), Inches(0.95),
             text, size=32, bold=True, color=TEAL_DARK, font=FONT_DISPLAY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.45), SW - Inches(1.2), Inches(0.5),
                 subtitle, size=16, color=GREY_SUB)


def dark_background(slide):
    add_rect(slide, 0, 0, SW, SH, TEAL_DARK)
    # accent circle-ish (rectangle band) for visual interest
    add_rect(slide, 0, 0, SW, Inches(0.15), GOLD)


def add_card(slide, x, y, w, h, *, fill=WHITE, border_right=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.05
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = RGBColor(0xe5, 0xdc, 0xc2)
    shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    if border_right:
        add_rect(slide, x + w - Emu(50000), y, Emu(50000), h, GOLD)
    return shp


# =========================================================
# Slide 1 — Cover
# =========================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, TEAL_DARK)
add_rect(s, 0, 0, SW, Inches(0.15), GOLD)
# gold circle-ish corner accent
acc = s.shapes.add_shape(MSO_SHAPE.OVAL, SW - Inches(4), -Inches(3), Inches(6), Inches(6))
acc.fill.solid(); acc.fill.fore_color.rgb = TEAL
acc.line.fill.background(); acc.shadow.inherit = False

add_text(s, Inches(1), Inches(1.8), SW - Inches(2), Inches(0.6),
         "يسرنا أن نعلن عن", size=24, color=CREAM, font=FONT_DISPLAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.4), SW - Inches(2), Inches(1.8),
         "معهد الهدي النبوي\nللعلوم الشرعية",
         size=58, bold=True, color=WHITE, font=FONT_DISPLAY, align=PP_ALIGN.CENTER)
# divider
add_rect(s, SW/2 - Inches(2), Inches(4.55), Inches(4), Emu(12000), GOLD)
add_text(s, Inches(1), Inches(4.7), SW - Inches(2), Inches(0.9),
         "بدء التسجيل على الدورة الأولى لعام ٢٠٢٦",
         size=34, color=GOLD, font=FONT_DISPLAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.7), SW - Inches(2), Inches(0.6),
         "برنامج تأهيل شرعي لغير المتفرغين",
         size=18, color=CREAM, align=PP_ALIGN.CENTER)
add_brand_strip(s); add_slide_number(s, 1, light=True)


# =========================================================
# Slide 2 — Vision / Mission / Goals
# =========================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, CREAM)
add_title(s, "الرؤيا والرسالة والأهداف", "من النظام الداخلي للمعهد")

# top row: 2 cards
row_y = Inches(2.1); row_h = Inches(1.85)
col_w = (SW - Inches(1.5)) / 2
gap = Inches(0.3)
card_w = col_w - gap/2

def vmg_card(slide, x, y, w, h, title, body):
    add_card(slide, x, y, w, h)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), w - Inches(0.4), Inches(0.55),
             title, size=22, bold=True, color=TEAL_DARK, font=FONT_DISPLAY)
    add_text(slide, x + Inches(0.2), y + Inches(0.7), w - Inches(0.4), h - Inches(0.8),
             body, size=14, color=GREY_TXT)

vmg_card(s, Inches(0.5) + col_w + gap/2, row_y, card_w, row_h,
         "الرؤيا",
         "الإسهام في ترسيخ الفهم المنضبط لمبادئ علوم الشريعة وأساسياتها، بما يستثمر تراث العلماء ويستجيب لمتطلبات العصر.")
vmg_card(s, Inches(0.5), row_y, card_w, row_h,
         "الرسالة",
         "تأهيل طلاب علم ودعاة يجمعون بين العلم والحكمة والتزكية، لبناء الوطن من خلال خطاب إسلامي يخاطب العقول ويجذب القلوب، وفق برامج تعليمية متكاملة لغير المتفرغين.")

# bottom wide card: goals
gy = Inches(4.1); gh = Inches(3.0); gw = SW - Inches(1.0)
add_card(s, Inches(0.5), gy, gw, gh)
add_text(s, Inches(0.7), gy + Inches(0.1), gw - Inches(0.4), Inches(0.5),
         "الأهداف", size=22, bold=True, color=TEAL_DARK, font=FONT_DISPLAY)
goals = [
    "تمكين الأصول العلمية في العلوم الشرعية لدى الدعاة إلى الله تعالى.",
    "الجمع بين تمكين العلم وتهذيب السلوك في الداعية إلى الله تعالى.",
    "تخريج دعاة يجمعون بين الصلاح الشخصي والإصلاح الاجتماعي.",
    "نشر العلوم الشرعية التي تصلح المجتمع وتعين على معالجة مشكلاته.",
    "تمكين أسس الخطاب الديني المعتدل الذي يعبّر عن روح الشريعة وسماحتها.",
]
add_text(s, Inches(0.7), gy + Inches(0.7), gw - Inches(0.4), gh - Inches(0.8),
         "\n".join("✦  " + g for g in goals), size=15, color=GREY_TXT)
add_brand_strip(s); add_slide_number(s, 2)


# =========================================================
# Slide 3 — Perks
# =========================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, CREAM)
add_title(s, "ميزات الدراسة في المعهد", "لماذا تختار معهد الهدي النبوي")

perks = [
    ("منهج علمي متين", "٢٤ مقرراً موزعة على ثلاثة فصول دراسية"),
    ("مرونة تامة", "مناسب لغير المتفرغين من العاملين والجامعيين"),
    ("اعتماد رسمي", "شهادة صادرة عن وزارة الأوقاف"),
    ("مجاني بالكامل", "لا رسوم ولا أقساط طوال مدة الدراسة"),
]
pw = (SW - Inches(1.5)) / 2
ph = Inches(2.2)
py0 = Inches(2.3)
for i, (t, b) in enumerate(perks):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + (1 - col) * (pw + Inches(0.5))  # RTL: first item on right
    y = py0 + row * (ph + Inches(0.3))
    add_card(s, x, y, pw, ph)
    add_text(s, x + Inches(0.3), y + Inches(0.25), pw - Inches(0.6), Inches(0.7),
             t, size=26, bold=True, color=TEAL_DARK, font=FONT_DISPLAY)
    add_text(s, x + Inches(0.3), y + Inches(1.05), pw - Inches(0.6), ph - Inches(1.2),
             b, size=16, color=GREY_SUB)
add_brand_strip(s); add_slide_number(s, 3)


# =========================================================
# Slide 4 — Study system
# =========================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, CREAM)
add_title(s, "نظام الدراسة والامتحانات", "نظرة عامة على المدة والتقييم")

stats = [("٣", "فصول دراسية"), ("٢١", "شهراً إجمالاً"),
         ("٢٤", "مقرراً"), ("٦٠٪", "درجة النجاح")]
sw_card = (SW - Inches(1.5)) / 4 - Inches(0.1)
sy = Inches(2.1)
for i, (n, l) in enumerate(stats):
    x = SW - Inches(0.5) - (i + 1) * (sw_card + Inches(0.15)) + Inches(0.15)
    add_card(s, x, sy, sw_card, Inches(1.5))
    add_text(s, x, sy + Inches(0.15), sw_card, Inches(0.8),
             n, size=44, bold=True, color=TEAL, font=FONT_DISPLAY, align=PP_ALIGN.CENTER)
    add_text(s, x, sy + Inches(1.0), sw_card, Inches(0.4),
             l, size=14, color=GREY_SUB, align=PP_ALIGN.CENTER)

rx = Inches(0.5); ry = Inches(3.9); rw = SW - Inches(1.0); rh = Inches(3.1)
add_card(s, rx, ry, rw, rh)
add_text(s, rx + Inches(0.3), ry + Inches(0.15), rw - Inches(0.5), Inches(0.5),
         "أبرز ضوابط الدراسة", size=20, bold=True, color=TEAL_DARK, font=FONT_DISPLAY)
rules = [
    "كل فصل يمتد ستة أشهر يتبعها شهر للعطلة والامتحانات.",
    "توزيع الدرجات: ٧٠٪ للامتحان النظري و٣٠٪ لأعمال السنة (الحضور والوظائف).",
    "نسبة حضور لا تقل عن ٧٠٪ شرط لدخول الامتحان النهائي.",
    "ينتقل الطالب إلى الفصل التالي إذا لم يرسب في أكثر من ثلاث مواد تراكمية.",
    "تقام دورة امتحانية بعد كل فصل، ودورة تكميلية عند اللزوم.",
]
add_text(s, rx + Inches(0.3), ry + Inches(0.75), rw - Inches(0.5), rh - Inches(0.9),
         "\n".join("✦  " + r for r in rules), size=14, color=GREY_TXT)
add_brand_strip(s); add_slide_number(s, 4)


# =========================================================
# Curriculum slides helper
# =========================================================
def curriculum_slide(n, term_title, subtitle, rows):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, CREAM)
    add_title(s, term_title, subtitle)

    tx = Inches(0.5); ty = Inches(2.1); tw = SW - Inches(1.0)
    nrows = len(rows) + 1
    th = Inches(4.9)
    # table
    table_shape = s.shapes.add_table(nrows, 3, tx, ty, tw, th)
    table = table_shape.table
    # column widths (RTL visual right→left): col0=subject, col1=hours, col2=book
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = tw - Inches(5.0)

    headers = ["المقرر", "الساعات", "الكتاب المقترح"]
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_DARK
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; set_rtl(p)
        r = p.add_run(); r.text = h
        r.font.name = FONT_DISPLAY; r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = CREAM

    for ri, (subj, hrs, book) in enumerate(rows, start=1):
        bg = RGBColor(0xfb, 0xf6, 0xe9) if ri % 2 == 0 else WHITE
        for ci, (val, col, bold, font) in enumerate([
            (subj, TEAL_DARK, True, FONT_DISPLAY),
            (hrs, RGBColor(0xb8, 0x8c, 0x1e), True, FONT_BODY),
            (book, GREY_TXT, False, FONT_BODY),
        ]):
            cell = table.cell(ri, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            tf = cell.text_frame; tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; set_rtl(p)
            r = p.add_run(); r.text = val
            r.font.name = font
            r.font.size = Pt(11 if ci == 2 else 12)
            r.font.bold = bold
            r.font.color.rgb = col
    add_brand_strip(s); add_slide_number(s, n)


term1 = [
    ("قرآن كريم ١", "ساعة / أسبوع", "حفظ جزئي عمّ وتبارك مع تطبيق أحكام التجويد الأساسية"),
    ("علوم القرآن والتفسير ١", "ساعة / أسبوع", "علوم القرآن — د. نور الدين عتر، والتفسير الوجيز للواحدي"),
    ("الحديث الشريف وعلومه ١", "ساعة / أسبوع", "شرح المنظومة البيقونية — الشيخ عبد الله سراج الدين + حفظ ١٠٠ حديث من رياض الصالحين"),
    ("فقه ١", "ساعتان / أسبوع", "فقه العبادات — للسلقيني"),
    ("نحو ١", "ساعة / أسبوع", "النحو الواضح (الجزء الأول) — علي الجارم"),
    ("عقيدة ١", "ساعة / أسبوع", "أوجز البيان في شرح أركان الإيمان — محمد صلاح تقوى"),
    ("السيرة النبوية ١", "ساعة / أسبوع", "روضة الأنوار — للمباركفوري"),
    ("تزكية وأخلاق ١", "ساعة / أسبوع", "أيها الولد — للغزالي، وتعليم المتعلم — للزرنوجي"),
]
term2 = [
    ("قرآن كريم ٢", "ساعة / أسبوع", "حفظ السور المسنونة (الكهف، يس، الواقعة، الدخان، السجدة)"),
    ("تفسير ٢", "ساعة / أسبوع", "تفسير السور المسنونة من صفوة التفاسير — للصابوني"),
    ("الحديث الشريف ٢", "ساعة / أسبوع", "إعلام الأنام في شرح بلوغ المرام — د. نور الدين عتر + حفظ ١٠٠ حديث"),
    ("فقه ٢", "ساعتان / أسبوع", "المفصل في الفقه الحنفي — الأموال والمعاملات، محمد ماجد عتر"),
    ("نحو ٢", "ساعة / أسبوع", "النحو الواضح (الجزء الثاني) — علي الجارم"),
    ("سيرة ٢", "ساعة / أسبوع", "إتمام الوفاء في سيرة الخلفاء — للخضري"),
    ("تزكية وأخلاق ٢", "ساعة / أسبوع", "موعظة القلوب — د. إبراهيم شاشو"),
    ("الخطابة وفن الدعوة", "ساعة / أسبوع", "الدعوة والداعية — د. نور الدين عتر، ومقرر الدعوة والخطابة (وزارة الأوقاف)"),
]
term3 = [
    ("قرآن كريم ٣", "ساعة / أسبوع", "حفظ سورة البقرة"),
    ("تفسير ٣", "ساعة / أسبوع", "تفسير سورة البقرة من صفوة التفاسير — للصابوني"),
    ("الحديث الشريف ٣", "ساعة / أسبوع", "إعلام الأنام في شرح بلوغ المرام — د. نور الدين عتر + حفظ ١٠٠ حديث"),
    ("فقه ٣", "ساعتان / أسبوع", "الأحوال الشخصية — قدري باشا"),
    ("نحو ٣", "ساعة / أسبوع", "النحو الواضح (الجزء الثالث)، والمختار في قواعد الإملاء — د. عبد البديع النيرباني"),
    ("سيرة ٣", "ساعة / أسبوع", "تاريخ الإسلام — للسرجاني (الأموية، العباسية، العثمانية)"),
    ("تزكية وأخلاق ٣", "ساعة / أسبوع", "موعظة القلوب — د. إبراهيم شاشو"),
    ("فِرَق ومذاهب", "ساعة / أسبوع", "الفرق الإسلامية — محمد أحمد خطيب"),
]

curriculum_slide(5, "الخطة الدراسية — الفصل الأول", "ثمانية مقررات أسبوعية", term1)
curriculum_slide(6, "الخطة الدراسية — الفصل الثاني", "توسّع في الفقه والتفسير ومدخل إلى الخطابة", term2)
curriculum_slide(7, "الخطة الدراسية — الفصل الثالث", "تعمّق في البقرة والأحوال الشخصية والتاريخ الإسلامي", term3)


# =========================================================
# Slide 8 — Requirements + Documents
# =========================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, CREAM)
add_title(s, "شروط القبول والوثائق المطلوبة", "يرجى إحضار الوثائق كاملة عند التسجيل")

cw = (SW - Inches(1.5)) / 2
cy = Inches(2.2); ch = Inches(4.4)
# right card — requirements (visually first in RTL)
rx = Inches(0.5) + cw + Inches(0.5)
lx = Inches(0.5)
for x, title, items in [
    (rx, "شروط القبول", [
        "أن يكون العمر ١٨ عاماً فما فوق",
        "شهادة التعليم الأساسي كحد أدنى",
        "اجتياز اختبار القبول الكتابي والشفهي",
        "التعهد بالالتزام بالدوام والحضور طوال مدة الدراسة",
    ]),
    (lx, "الوثائق المطلوبة", [
        "صورة عن الهوية الشخصية (عدد ٢)",
        "صورة شخصية حديثة ملوّنة (عدد ٢)",
        "صورة مصدّقة عن المؤهل العلمي",
        "تعبئة نموذج تعهد الالتزام",
    ])
]:
    add_card(s, x, cy, cw, ch)
    add_text(s, x + Inches(0.3), cy + Inches(0.2), cw - Inches(0.5), Inches(0.7),
             title, size=26, bold=True, color=TEAL_DARK, font=FONT_DISPLAY)
    add_text(s, x + Inches(0.3), cy + Inches(1.0), cw - Inches(0.5), ch - Inches(1.2),
             "\n".join("✦  " + it for it in items), size=16, color=GREY_TXT)
add_brand_strip(s); add_slide_number(s, 8)


# =========================================================
# Slide 9 — CTA
# =========================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, TEAL_DARK)
add_rect(s, 0, 0, SW, Inches(0.15), GOLD)
acc = s.shapes.add_shape(MSO_SHAPE.OVAL, -Inches(3), SH - Inches(3), Inches(6), Inches(6))
acc.fill.solid(); acc.fill.fore_color.rgb = TEAL
acc.line.fill.background(); acc.shadow.inherit = False

add_text(s, Inches(1), Inches(1.0), SW - Inches(2), Inches(0.6),
         "للتسجيل في الدورة الأولى ٢٠٢٦",
         size=24, color=CREAM, font=FONT_DISPLAY, align=PP_ALIGN.CENTER)

# dates pill
pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          SW/2 - Inches(4), Inches(1.9), Inches(8), Inches(1.1))
pill.adjustments[0] = 0.3
pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor(0x1c, 0x55, 0x50)
pill.line.color.rgb = GOLD; pill.line.width = Pt(1)
pill.shadow.inherit = False
tb = s.shapes.add_textbox(SW/2 - Inches(4), Inches(2.0), Inches(8), Inches(0.9))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; set_rtl(p)
for txt, col, bold in [
    ("آخر موعد لاستقبال الطلبات: ", CREAM, False),
    ("١٩ نيسان", GOLD, True),
]:
    r = p.add_run(); r.text = txt
    r.font.name = FONT_DISPLAY; r.font.size = Pt(30); r.font.bold = bold
    r.font.color.rgb = col

add_text(s, Inches(2), Inches(3.3), SW - Inches(4), Inches(0.8),
         "التسجيل حضورياً في مقر المعهد لتسليم الوثائق ومقابلة القبول",
         size=20, color=CREAM, align=PP_ALIGN.CENTER)

# place card
px = SW/2 - Inches(3.5); py = Inches(4.4); pw = Inches(7); ph = Inches(2.0)
place = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, py, pw, ph)
place.adjustments[0] = 0.08
place.fill.solid(); place.fill.fore_color.rgb = RGBColor(0x10, 0x42, 0x3c)
place.line.color.rgb = GOLD; place.line.width = Pt(1)
place.shadow.inherit = False
add_rect(s, px + pw - Emu(60000), py, Emu(60000), ph, GOLD)
add_text(s, px + Inches(0.3), py + Inches(0.2), pw - Inches(0.5), Inches(0.5),
         "مقر التسجيل", size=16, bold=True, color=GOLD, font=FONT_DISPLAY, align=PP_ALIGN.CENTER)
add_text(s, px + Inches(0.3), py + Inches(0.75), pw - Inches(0.5), ph - Inches(0.9),
         "معهد القرآن الكريم — جامع خالد بن الوليد\nالحمدانية",
         size=22, color=CREAM, font=FONT_DISPLAY, align=PP_ALIGN.CENTER)

add_brand_strip(s); add_slide_number(s, 9, light=True)


out = "/Users/alisafaya/Desktop/معهد-الهدي-النبوي/infographic/presentation-draft.pptx"
prs.save(out)
print(f"Wrote {out}")
